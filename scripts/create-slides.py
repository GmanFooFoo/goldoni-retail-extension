#!/usr/bin/env python3
"""
Goldoni Retail — PPTX-Generator

Liest Slide-Inhalte aus slides-content.yaml, erzeugt goldoni-retail-ueberblick.pptx.
Farbschema: Dunkelbraun / Creme / Gold (Restaurant-Ästhetik).

Usage:
    python3 scripts/create-slides.py
    python3 scripts/create-slides.py --input pfad/zu/content.yaml --output pfad/zu/output.pptx
"""

import argparse
import os
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Farben ──────────────────────────────────────────────────
BG_DARK = RGBColor(0x2C, 0x24, 0x1E)
BG_CREAM = RGBColor(0xF5, 0xF0, 0xE6)
TEXT_CREAM = RGBColor(0xF5, 0xF0, 0xE6)
TEXT_DARK = RGBColor(0x2C, 0x24, 0x1E)
GOLD = RGBColor(0xC8, 0xA2, 0x5C)
GREEN = RGBColor(0x6B, 0x8E, 0x4E)
RED_SOFT = RGBColor(0xA0, 0x4A, 0x4A)
GRAY = RGBColor(0xD0, 0xC8, 0xBC)
CELL_BG = RGBColor(0x3A, 0x32, 0x2A)
FONT = "Georgia"


# ── Hilfs-Funktionen ───────────────────────────────────────
def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK


def gold_line(slide, top):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), top, Inches(11.7), Pt(2))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()


def tbox(slide, l, t, w, h, text, size=18, color=TEXT_CREAM,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return tf


def mtext(slide, l, t, w, h, lines, size=18, color=TEXT_CREAM, spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        if isinstance(ln, dict):
            txt = ln.get("text", "")
            bld = ln.get("bold", False)
        else:
            txt, bld = str(ln), False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bld
        p.font.name = FONT
        p.space_after = Pt(size * spacing * 0.4)
    return tf


def add_table(slide, l, t, w, rows, font_size=14,
              hdr_bg=GOLD, hdr_txt=TEXT_DARK, body_txt=TEXT_CREAM):
    n_rows, n_cols = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(l), Inches(t), Inches(w),
        Inches(0.45 * n_rows)).table
    col_w = Inches(w / n_cols)
    for c in range(n_cols):
        tbl.columns[c].width = col_w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.font.name = FONT
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = hdr_txt
                else:
                    para.font.color.rgb = body_txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = hdr_bg if r == 0 else CELL_BG
    return tbl


def headline(slide, text):
    tbox(slide, 0.8, 0.5, 11.7, 0.8, text, size=36, color=GOLD, bold=True)
    gold_line(slide, Inches(1.3))


def footer(slide, text, color=GOLD, bold=True):
    tbox(slide, 0.8, 6.2, 11.7, 0.8, text, size=18, color=color, bold=bold)


# ── Slide-Builder je Typ ───────────────────────────────────
def build_title(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    gold_line(s, Inches(3.0))
    tbox(s, 0.8, 1.5, 11.7, 1.5, d["headline"],
         size=54, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    tbox(s, 0.8, 3.3, 11.7, 1, d.get("subline", ""),
         size=24, align=PP_ALIGN.CENTER)
    if d.get("footer"):
        tbox(s, 0.8, 5.5, 11.7, 0.5, d["footer"],
             size=16, color=GRAY, align=PP_ALIGN.CENTER)


def build_text(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    mtext(s, 0.8, 1.8, 11.7, 4.5, d["lines"], size=22)


def build_text_and_table(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    mtext(s, 0.8, 1.8, 5.5, 2, d.get("intro", []), size=20)
    if d.get("table_label"):
        tbox(s, 0.8, 3.8, 5.5, 0.5, d["table_label"], size=20, bold=True)
    add_table(s, 0.8, 4.5, 11.7, d["table"])
    if d.get("footer"):
        footer(s, d["footer"], color=GREEN)


def build_checklist(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    for i, item in enumerate(d["items"]):
        y = 1.8 + i * 0.7
        tbox(s, 1.0, y, 0.5, 0.5, "✓", size=22, color=GREEN, bold=True)
        tbox(s, 1.6, y, 10.5, 0.5, item, size=20)
    if d.get("footer"):
        footer(s, d["footer"])


def build_table(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    if d.get("subtitle"):
        tbox(s, 0.8, 1.6, 11.7, 0.6, d["subtitle"], size=22, bold=True)
        add_table(s, 0.8, 2.4, 11.7, d["table"], font_size=16)
    else:
        add_table(s, 0.8, 1.8, 11.7, d["table"], font_size=16)
    if d.get("footer"):
        tbox(s, 0.8, 5.0, 11.7, 1.5, d["footer"], size=18)


def build_table_and_sidebar(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    if d.get("subtitle"):
        tbox(s, 0.8, 1.6, 11.7, 0.6, d["subtitle"], size=26, bold=True)
    add_table(s, 0.8, 2.4, 7, d["table"], font_size=15)
    # Sidebar
    sb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(8.5), Inches(2.4), Inches(4.3), Inches(3.5))
    sb.fill.solid()
    sb.fill.fore_color.rgb = CELL_BG
    sb.line.fill.background()
    mtext(s, 8.8, 2.6, 3.8, 3, d.get("sidebar", []), size=16)


def build_table_and_text(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    if d.get("subtitle"):
        tbox(s, 0.8, 1.6, 11.7, 0.6, d["subtitle"], size=22, bold=True)
        add_table(s, 0.8, 2.4, 11.7, d["table"], font_size=15)
        mtext(s, 0.8, 5.0, 11.7, 2, d.get("lines", []), size=18)
    else:
        add_table(s, 0.8, 1.8, 11.7, d["table"], font_size=18)
        mtext(s, 0.8, 3.8, 11.7, 2.5, d.get("lines", []), size=20)


def build_numbered_list(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    for i, item in enumerate(d["items"]):
        y = 1.7 + i * 0.6
        tbox(s, 0.8, y, 0.5, 0.5, f"{i+1}.", size=18, color=GOLD, bold=True)
        tbox(s, 1.4, y, 4.5, 0.5, item["action"], size=18, bold=True)
        tbox(s, 6.0, y, 6.5, 0.5, item.get("detail", ""), size=16, color=GRAY)
    if d.get("footer"):
        footer(s, d["footer"])


def build_bullet_list(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    for i, item in enumerate(d["items"]):
        y = 1.7 + i * 0.6
        tbox(s, 1.0, y, 0.5, 0.5, "—", size=18, color=GOLD)
        tbox(s, 1.5, y, 10.5, 0.5, item, size=18)
    if d.get("footer"):
        tbox(s, 0.8, 6.2, 11.7, 0.8, d["footer"], size=18, color=GRAY)


def build_risks(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    for i, item in enumerate(d["items"]):
        y = 1.7 + i * 1.0
        tbox(s, 0.8, y, 0.4, 0.4, str(i+1), size=16, color=RED_SOFT, bold=True)
        tbox(s, 1.3, y, 3, 0.4, item["title"], size=17, bold=True)
        tbox(s, 1.3, y + 0.35, 10.5, 0.6, item["detail"], size=15, color=GRAY)


def build_closing(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    headline(s, d["headline"])
    mtext(s, 0.8, 1.8, 11.7, 2, d.get("intro", []), size=24)
    # Zitat-Box
    sb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(1.5), Inches(3.8), Inches(10.3), Inches(2))
    sb.fill.solid()
    sb.fill.fore_color.rgb = CELL_BG
    sb.line.fill.background()
    quote_lines = [{"text": q, "bold": False} for q in d.get("quote", [])]
    mtext(s, 1.8, 4.0, 9.7, 1.8, quote_lines, size=20, color=BG_CREAM)
    if d.get("footer"):
        footer(s, d["footer"])


# ── Dispatch ───────────────────────────────────────────────
BUILDERS = {
    "title": build_title,
    "text": build_text,
    "text-and-table": build_text_and_table,
    "checklist": build_checklist,
    "table": build_table,
    "table-and-sidebar": build_table_and_sidebar,
    "table-and-text": build_table_and_text,
    "numbered-list": build_numbered_list,
    "bullet-list": build_bullet_list,
    "risks": build_risks,
    "closing": build_closing,
}


def main():
    parser = argparse.ArgumentParser(description="Goldoni Retail PPTX-Generator")
    script_dir = os.path.dirname(os.path.abspath(__file__))
    repo_root = os.path.dirname(script_dir)
    default_input = os.path.join(repo_root,
                                 "docs/silvio-derivatives/slides-content.yaml")
    default_output = os.path.join(repo_root,
                                  "docs/silvio-derivatives/goldoni-retail-ueberblick.pptx")

    parser.add_argument("--input", "-i", default=default_input,
                        help="YAML-Datei mit Slide-Inhalten")
    parser.add_argument("--output", "-o", default=default_output,
                        help="Ausgabe-PPTX")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in data["slides"]:
        slide_type = slide_data["type"]
        builder = BUILDERS.get(slide_type)
        if not builder:
            print(f"Unbekannter Slide-Typ: {slide_type} — übersprungen.")
            continue
        builder(prs, slide_data)

    prs.save(args.output)
    print(f"Gespeichert: {args.output}")
    print(f"{len(data['slides'])} Slides generiert.")


if __name__ == "__main__":
    main()
